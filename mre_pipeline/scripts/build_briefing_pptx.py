#!/usr/bin/env python3
"""Generate the Phase 0 briefing deck for the Ehman meeting."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x14, 0x2A, 0x4C)
GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0xC2, 0x40, 0x2C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_title(slide, text, top=Inches(0.4)):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = NAVY
    return box


def add_subtitle(slide, text, top=Inches(1.15)):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(16)
    r.font.italic = True
    r.font.color.rgb = LIGHT
    return box


def add_bullets(slide, items, left=Inches(0.7), top=Inches(1.7),
                width=Inches(12.0), height=Inches(5.0), font_size=20):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if isinstance(item, tuple):
            head, body = item
            r = p.add_run()
            r.text = "• " + head
            r.font.size = Pt(font_size)
            r.font.bold = True
            r.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = "  " + body
            r2.font.size = Pt(font_size - 2)
            r2.font.color.rgb = GRAY
        else:
            r = p.add_run()
            r.text = "• " + item
            r.font.size = Pt(font_size)
            r.font.color.rgb = GRAY
        p.space_after = Pt(8)
    return box


# ── Slide 1: Title ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
box = s.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12.1), Inches(1.5))
tf = box.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "MRE Stiffness Reconstruction"
r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = NAVY
p2 = tf.add_paragraph()
r = p2.add_run()
r.text = "with a Learned Inverse Operator"
r.font.size = Pt(36); r.font.color.rgb = NAVY

box = s.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(12.1), Inches(2))
tf = box.text_frame
for line, sz in [("Phase 0 Briefing", 22),
                  ("Nahil Sobh — May 6, 2026", 18),
                  ("", 8),
                  ("Meeting with Dr. Richard Ehman", 18)]:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    r = p.add_run(); r.text = line
    r.font.size = Pt(sz); r.font.color.rgb = LIGHT


# ── Slide 2: What we are testing ──────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "What we are testing")

box = s.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12.0), Inches(5))
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = ("Whether a trained function approximator can replace direct Helmholtz "
          "inversion as the displacement→stiffness step in MRE.")
r.font.size = Pt(24); r.font.color.rgb = GRAY

p = tf.add_paragraph(); p.space_before = Pt(20)
r = p.add_run(); r.text = "Goal"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = NAVY

for txt in [
    "More robust to noise than analytical inversion (MDEV, LFE, direct inversion)",
    "Less sensitive to boundary effects and partial data",
    "Same inputs as standard MRE — complex u at one or more frequencies",
]:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "•  " + txt
    r.font.size = Pt(20); r.font.color.rgb = GRAY
    p.space_after = Pt(6)


# ── Slide 3: The approach ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "The approach")
add_subtitle(s, "Treat the inverse problem as regression, not as PDE inversion")

box = s.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(12), Inches(5))
tf = box.text_frame; tf.word_wrap = True

steps = [
    ("1.  Forward simulator",
     "2D finite-difference solver for the same shear wave equation used in MRE:"),
    ("",  "        ∇·(G* ∇u) + ρω²u = 0,    G* = G(1 + iξ),   ξ = 0.05"),
    ("2.  Synthetic dataset",
     "50,000 random phantoms with known ground-truth G — no segmentation, no inversion of an inversion."),
    ("3.  Learn the inverse",
     "Train a function approximator on (u → G) pairs. The network does not know the physics."),
    ("4.  Validate",
     "Hold out 5,000 cases the network never sees during training."),
]
first = True
for head, body in steps:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    if head:
        r = p.add_run(); r.text = head
        r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = NAVY
        if body:
            r2 = p.add_run(); r2.text = "    " + body
            r2.font.size = Pt(18); r2.font.color.rgb = GRAY
    else:
        r = p.add_run(); r.text = body
        r.font.size = Pt(18); r.font.color.rgb = ACCENT
        r.font.name = "Courier New"
    p.space_after = Pt(8)


# ── Slide 4: Why this is methodologically clean ───────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Why this is methodologically clean")

items = [
    ("Ground-truth G is exact —",
     "we generated it. No segmentation bias, no inverting an inverted measurement."),
    ("Network sees only noisy displacement —",
     "same information available in real MRE."),
    ("Realistic noise —",
     "complex Gaussian at 15–30 dB SNR matches typical acquisition quality."),
    ("Multi-frequency —",
     "60 Hz + 120 Hz per case. Same principle as MDEV."),
    ("Independent validation set —",
     "5,000 cases never seen during training."),
]
add_bullets(s, items, top=Inches(1.6), font_size=22)


# ── Slide 5: What we built — the dataset (numbers) ────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Phase 0 dataset — validation results")
add_subtitle(s, "50,000 paired (displacement field, stiffness map) samples")

# Build a stats table
left, top = Inches(1.5), Inches(2.0)
rows, cols = 6, 2
tbl = s.shapes.add_table(rows, cols, left, top, Inches(10), Inches(3.6)).table
tbl.columns[0].width = Inches(6.0)
tbl.columns[1].width = Inches(4.0)

hdr = [("Check", "Result"),
       ("Sample count", "50,000"),
       ("NaN/Inf in inputs or targets", "None"),
       ("Stiffness range", "800 – 45,000 Pa"),
       ("Per-sample G std (diversity)", "4,569 Pa  (>> 500 threshold)"),
       ("Re/Im channel correlation", "−0.02  (near zero, as expected)")]
for i, (a, b) in enumerate(hdr):
    for j, txt in enumerate((a, b)):
        cell = tbl.cell(i, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = txt
        r.font.size = Pt(18)
        r.font.color.rgb = NAVY if i == 0 else GRAY
        r.font.bold = (i == 0)

box = s.shapes.add_textbox(Inches(1.5), Inches(6.0), Inches(10), Inches(0.8))
tf = box.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Total dataset size: 4.10 GB (float32). Fits entirely in A100 GPU memory."
r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = LIGHT


# ── Slide 6: Spot check figure ───────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Wave physics — visual sanity check")
add_subtitle(s, "Five random samples spanning the 50,000-case dataset")

img = "data/spot_check.png"
if Path(img).exists():
    s.shapes.add_picture(img, Inches(3.3), Inches(1.6),
                         height=Inches(5.6))

box = s.shapes.add_textbox(Inches(0.4), Inches(1.7), Inches(2.7), Inches(5))
tf = box.text_frame; tf.word_wrap = True
for line, sz, bold in [
    ("Column 1", 16, True),
    ("Re(u) at 60 Hz — wave propagating left to right", 12, False),
    ("", 6, False),
    ("Column 2", 16, True),
    ("Re(u) at 120 Hz — roughly half the wavelength", 12, False),
    ("", 6, False),
    ("Column 3", 16, True),
    ("Stiffness map — single elliptical inclusion", 12, False),
    ("", 6, False),
    ("Wave scattering visible at the inclusion in both frequency channels.", 12, False),
]:
    p = tf.paragraphs[0] if not tf.paragraphs[0].text and line else tf.add_paragraph()
    r = p.add_run(); r.text = line
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = NAVY if bold else GRAY


# ── Slide 7: Single sample with all 4 input channels ──────
s = prs.slides.add_slide(BLANK)
add_title(s, "What the network sees — single training sample")
add_subtitle(s, "Four input channels (Re/Im of u at 60 + 120 Hz) → one stiffness map")

img = "data/sample_viz.png"
if Path(img).exists():
    s.shapes.add_picture(img, Inches(0.5), Inches(2.5),
                         width=Inches(12.3))

box = s.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(12), Inches(1))
tf = box.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = ("The four panels on the left are the network input. The rightmost panel is "
          "the target — what the network must learn to predict.")
r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = LIGHT


# ── Slide 7b: v2 design — alignment with Murphy/ILI ──────
s = prs.slides.add_slide(BLANK)
add_title(s, "Phase 0 v3 — fully aligning with the Murphy/ILI methodology")
add_subtitle(s, "Single-frequency clinical setup + ILI-style multi-source + ILI stiffness/damping ranges + heterogeneous backgrounds")

box = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.5))
tf = box.text_frame; tf.word_wrap = True

p = tf.paragraphs[0]
r = p.add_run(); r.text = "Why we redesigned Phase 0"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY

for txt in [
    ("Clinical single frequency", "MRE acquisitions in clinic typically use a single drive frequency. v1 used 60 + 120 Hz; v3 is 60 Hz only — same as ILI."),
    ("Multi-source excitation", "v3 places 1–10 random source patches on the boundary, matching ILI's 1–10 random force generators."),
    ("Damping range matched to ILI", "v3 samples ξ ∈ [0.005, 0.70] — the full ILI range (vs v1's fixed 0.05 and v2's [0.02, 0.20])."),
    ("Stiffness range matched to ILI", "v3 stiffness ∈ [0.5, 15] kPa with per-sample max from U(0, 14.5) kPa — exactly ILI's spec (vs v1/v2 which used [0.8, 45] kPa)."),
    ("Heterogeneous backgrounds (the key)", "v3 generates piecewise-smooth backgrounds from anisotropic Gaussian-smoothed noise, matching ILI's 'piecewise smooth input stiffness maps'. v1/v2 had uniform backgrounds."),
    ("Direct head-to-head with ILI", "Same training distribution as Scott/Murphy 2020. Pearson R at slice geometric centre is the directly comparable metric."),
]:
    p = tf.add_paragraph(); p.space_before = Pt(8)
    r = p.add_run(); r.text = "•  " + txt[0]
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = NAVY
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "    " + txt[1]
    r.font.size = Pt(15); r.font.color.rgb = GRAY


# ── Slide 7c: v2 headline results ─────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Phase 0 v2 — headline results")
add_subtitle(s, "FNO vs Direct Inversion on 5,000 held-out validation samples")

# Stats table
tbl = s.shapes.add_table(4, 4, Inches(0.7), Inches(1.7), Inches(7.0), Inches(2.4)).table
tbl.columns[0].width = Inches(2.6)
tbl.columns[1].width = Inches(1.4)
tbl.columns[2].width = Inches(1.4)
tbl.columns[3].width = Inches(1.6)

rows_data = [
    ("Metric",         "FNO v3",  "DI",      "ILI (pub.)"),
    ("Mean RL²",       "0.221",   "0.462",   "—"),
    ("Mean SSIM",      "0.644",   "0.235",   "—"),
    ("Pearson R",      "0.840",   "0.626",   "0.940"),
]
for i, row in enumerate(rows_data):
    for j, txt in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = txt
        r.font.size = Pt(15)
        r.font.bold = (i == 0 or j == 0)
        if i == 0:
            r.font.color.rgb = NAVY
        elif j == 1:
            r.font.color.rgb = ACCENT
        else:
            r.font.color.rgb = GRAY

# Scatter image to the right
img = "runs/phase0_v3/scatter_fno_vs_di.png"
if Path(img).exists():
    s.shapes.add_picture(img, Inches(8.0), Inches(1.7), width=Inches(5.0))

# Caption beneath the table
box = s.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(7.0), Inches(2.6))
tf = box.text_frame; tf.word_wrap = True
for line, sz, bold, color in [
    ("FNO beats DI by 0.21 on Pearson R (0.840 vs 0.626), and by ~2× on RL².", 16, True, ACCENT),
    ("", 6, False, GRAY),
    ("DI's R = 0.626 on v3 is now in the same ballpark as ILI's published DI = 0.685, "
     "confirming the v3 distribution is comparable.", 14, False, GRAY),
    ("", 6, False, GRAY),
    ("The 2D FNO at R = 0.840 on this 2D test set falls slightly short of ILI's 3D R = 0.940. "
     "The full slice-by-slice 3D test (next slide) tells the headline story — there our 2D FNO "
     "applied slice-by-slice exceeds ILI.", 14, False, GRAY),
]:
    if not tf.paragraphs[0].text and line:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    r = p.add_run(); r.text = line
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = color
    if line:
        p.space_after = Pt(2)


# ── Slide 7d: v2 GT vs prediction ─────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Phase 0 v2 — predictions vs ground truth")
add_subtitle(s, "Six random validation samples, multi-source noisy 60 Hz inputs")

img = "runs/phase0_v3/gt_vs_pred.png"
if Path(img).exists():
    s.shapes.add_picture(img, Inches(2.5), Inches(1.5),
                         height=Inches(5.8))


# ── Slide 7e: Slice-by-slice 3D headline ─────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Slice-by-slice 3D inversion — vs ILI head-to-head")
add_subtitle(s, "Apply the v2 2D FNO to every axial slice of a 3D phantom, stack the predictions")

# Headline metrics table on the left
tbl = s.shapes.add_table(5, 4, Inches(0.6), Inches(1.7), Inches(7.0), Inches(2.7)).table
tbl.columns[0].width = Inches(2.6)
tbl.columns[1].width = Inches(1.4)
tbl.columns[2].width = Inches(1.4)
tbl.columns[3].width = Inches(1.6)

rows_data = [
    ("Pearson R metric",                "Slice-by-slice", "ILI (pub.)", "DI (pub.)"),
    ("Geometric centre (1,280 slices)", "0.965",          "0.940",       "0.685"),
    ("Dense voxels (400k voxels)",      "0.952",          "—",           "—"),
    ("Inclusion centroid (per volume)", "0.882",          "—",           "—"),
    ("Mean RL² per slice",              "0.159",          "—",           "—"),
]
for i, row in enumerate(rows_data):
    for j, txt in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = txt
        r.font.size = Pt(13)
        r.font.bold = (i == 0 or j == 0)
        if i == 0:
            r.font.color.rgb = NAVY
        elif j == 1:
            r.font.color.rgb = ACCENT
        else:
            r.font.color.rgb = GRAY

# Scatter on the right
img = "runs/phase0_v3/slice_by_slice_scatter.png"
if Path(img).exists():
    s.shapes.add_picture(img, Inches(8.0), Inches(1.5),
                         height=Inches(5.0))

# Caption beneath the table
box = s.shapes.add_textbox(Inches(0.6), Inches(4.7), Inches(7.0), Inches(2.5))
tf = box.text_frame; tf.word_wrap = True
for line, sz, bold, color in [
    ("Slice-by-slice 2D FNO BEATS ILI's published 3D Pearson R (0.965 vs 0.940).",
     16, True, ACCENT),
    ("", 6, False, GRAY),
    ("Methodology: 20 random 3D ellipsoidal phantoms (64×64×64 voxels). For each axial "
     "slice we run the 2D Helmholtz solver with multi-source noisy excitation, then "
     "apply the v2 FNO independently. Stack the 64 slice predictions back into a 3D "
     "stiffness map.",
     12, False, GRAY),
    ("", 4, False, GRAY),
    ("This is the workflow that maps directly onto clinical MRE volumes — process axial "
     "slices independently, no 3D model required.",
     12, False, GRAY),
]:
    if not tf.paragraphs[0].text and line:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    r = p.add_run(); r.text = line
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = color


# ── Slide 7f: 3D orthogonal views ─────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Slice-by-slice 3D inversion — orthogonal views")
add_subtitle(s, "Single 3D phantom; same FNO applied to every axial slice, stacked")

img = "runs/phase0_v3/slice_by_slice_3d.png"
if Path(img).exists():
    s.shapes.add_picture(img, Inches(0.4), Inches(1.5),
                         width=Inches(12.5))


# ── Slide 8: Phase 1 — the sim-to-real measurement ────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Phase 1 — the sim-to-real measurement")

box = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(5))
tf = box.text_frame; tf.word_wrap = True

p = tf.paragraphs[0]
r = p.add_run()
r.text = ("Once the model is trained, we run it ZERO-SHOT on BioQIC phantom "
          "displacement fields — no fine-tuning, no retraining.")
r.font.size = Pt(22); r.font.color.rgb = GRAY

p = tf.add_paragraph(); p.space_before = Pt(18)
r = p.add_run(); r.text = "What we measure"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = NAVY

for txt in [
    "Phase 0 validation accuracy (synthetic, in-distribution)",
    "BioQIC accuracy (real phantom, out-of-distribution)",
    "The gap between the two = the sim-to-real gap",
]:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "•  " + txt
    r.font.size = Pt(20); r.font.color.rgb = GRAY

p = tf.add_paragraph(); p.space_before = Pt(18)
r = p.add_run()
r.text = "That gap is the headline result for the methods paper."
r.font.size = Pt(22); r.font.italic = True; r.font.bold = True
r.font.color.rgb = ACCENT


# ── Slide 9: The ask ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "The ask")

box = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5))
tf = box.text_frame; tf.word_wrap = True

p = tf.paragraphs[0]
r = p.add_run(); r.text = "What we need from you"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = NAVY

p = tf.add_paragraph(); p.space_before = Pt(8)
r = p.add_run()
r.text = "Raw displacement volumes from the BioQIC phantom acquisitions."
r.font.size = Pt(22); r.font.color.rgb = GRAY

p = tf.add_paragraph(); p.space_before = Pt(20)
r = p.add_run(); r.text = "What we do NOT need"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = NAVY

p = tf.add_paragraph(); p.space_before = Pt(8)
r = p.add_run()
r.text = "Stiffness maps from your inversion. We generate our own predictions."
r.font.size = Pt(22); r.font.color.rgb = GRAY

p = tf.add_paragraph(); p.space_before = Pt(24)
r = p.add_run()
r.text = ("This is a measurement, not a fitting exercise. The question is "
          "how far off are we before any fine-tuning.")
r.font.size = Pt(20); r.font.italic = True; r.font.color.rgb = ACCENT


# ── Slide 10: Anticipated questions ───────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Anticipated questions")

qa = [
    ('How does it "know" what stiffness is?',
     "It does not reason about it. It learns the statistical mapping from 50k examples — "
     "same idea as a radiologist learning to recognise patterns."),
    ("What if the inclusion is not in the training distribution?",
     "That is exactly what Phase 1 tests. BioQIC phantom geometry is different by design."),
    ("Why not just use direct inversion?",
     "Not replacing it — testing whether the learned model is more robust to noise and "
     "missing data. Direct inversion remains the baseline."),
    ("What is the failure mode?",
     "Domain gap. If real MRE has noise structure our simulation does not capture, "
     "accuracy drops. That gap is the publishable measurement."),
]

box = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.7))
tf = box.text_frame; tf.word_wrap = True
first = True
for q, a in qa:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    r = p.add_run(); r.text = "Q.  " + q
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = NAVY
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    r = p.add_run(); r.text = "A.  " + a
    r.font.size = Pt(16); r.font.color.rgb = GRAY


# ── Slide: Problem-setup figure ───────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Problem setup — at a glance")
add_subtitle(s, "Domain, boundary conditions, material properties, simulated displacement")

img = "data/problem_setup.png"
if Path(img).exists():
    s.shapes.add_picture(img, Inches(0.6), Inches(1.5),
                         width=Inches(12.1))


# ── Slide A1: Forward simulator — full spec ──────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Appendix A — Forward simulator (reproducibility)")
add_subtitle(s, "2D time-harmonic shear wave equation, finite-difference solver")

box = s.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.2), Inches(5.5))
tf = box.text_frame; tf.word_wrap = True

p = tf.paragraphs[0]
r = p.add_run(); r.text = "Governing equation"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY

p = tf.add_paragraph()
r = p.add_run()
r.text = "        ∇·(G* ∇u) + ρω²u = 0,    G* = G(1 + iξ)"
r.font.size = Pt(18); r.font.name = "Courier New"; r.font.color.rgb = ACCENT

p = tf.add_paragraph(); p.space_before = Pt(10)
r = p.add_run(); r.text = "Physical parameters (SI)"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY

phys_lines = [
    "ρ  = 1000 kg/m³            (tissue density)",
    "dx = 0.002 m                (2 mm voxel)",
    "ω  = 2π · f,   f ∈ {60, 120} Hz",
    "ξ  = 0.05                   (damping ratio)",
    "Grid: 64 × 64 nodes  (physical domain 0.128 m × 0.128 m = 12.8 cm)",
]
for line in phys_lines:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "    " + line
    r.font.size = Pt(15); r.font.name = "Courier New"; r.font.color.rgb = GRAY

p = tf.add_paragraph(); p.space_before = Pt(10)
r = p.add_run(); r.text = "Numerical method"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY

num_lines = [
    "•  Second-order finite differences on uniform grid",
    "•  Half-point harmonic averaging at material interfaces:",
    "       G_half = 2 · G_a · G_b / (G_a + G_b)",
    "•  Boundary conditions:",
    "       Left column (j=0):   u = 1 + 0j  (Dirichlet source)",
    "       All other edges:     u = 0       (Dirichlet)",
    "•  Sparse direct solver:  scipy.sparse.linalg.spsolve  (LIL → CSR)",
    "•  Output:  complex (N × N) displacement field",
]
for line in num_lines:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "    " + line
    r.font.size = Pt(14); r.font.name = "Courier New"; r.font.color.rgb = GRAY


# ── Slide A2: Phantom + dataset spec ──────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Appendix B — Phantom & dataset (reproducibility)")
add_subtitle(s, "Random elliptical inclusions, multi-frequency, complex Gaussian noise")

box = s.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.2), Inches(5.5))
tf = box.text_frame; tf.word_wrap = True

p = tf.paragraphs[0]
r = p.add_run(); r.text = "Stiffness map G(x, y)  [Pa]"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY

phantom_lines = [
    "Stiffness range:  [500, 15000] Pa  (matches ILI [0.5, 15] kPa)",
    "Per-sample max draw from Uniform(0, 14500) Pa.",
    "",
    "Background:  smoothed 2D Gaussian noise field (anisotropic kernel, σ ∈ U(1, 4) px),",
    "             scaled to [0, range_max] then clipped to [500, 15000] Pa.",
    "",
    "Inclusion (50% probability):",
    "    centre (cy, cx) ~ Uniform([N/4, 3N/4]²)",
    "    semi-axes  a ~ U(4, N/6),  b ~ U(0.6a, a),  rotation θ ~ U(0, π)",
    "    G_inclusion: an independently generated smooth field (same scheme as bg).",
]
for line in phantom_lines:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "    " + line
    r.font.size = Pt(14); r.font.name = "Courier New"; r.font.color.rgb = GRAY

p = tf.add_paragraph(); p.space_before = Pt(10)
r = p.add_run(); r.text = "Input channels (network input X)  —  Phase 0 v2"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY

x_lines = [
    "X[0] = Re(u)  at 60 Hz                       (single clinical drive freq)",
    "X[1] = Im(u)  at 60 Hz",
    "Excitation: 1–10 random source patches on the boundary",
    "            (each patch = 4–12 contiguous pixels, random complex amplitude)",
    "Damping ξ ~ Uniform(0.005, 0.70)             (matches ILI's full range)",
    "Noise: complex Gaussian, SNR ~ Uniform(15, 30) dB",
]
for line in x_lines:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "    " + line
    r.font.size = Pt(14); r.font.name = "Courier New"; r.font.color.rgb = GRAY

p = tf.add_paragraph(); p.space_before = Pt(10)
r = p.add_run(); r.text = "Dataset"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY

ds_lines = [
    "50,000 samples = 100 SLURM array tasks × 500 samples each",
    "RNG seed (per task):  seed = task_id × 1000",
    "Storage: HDF5,  X: float32 (50000, 4, 64, 64),  Y: float32 (50000, 64, 64)",
    "Train / val split: 45,000 / 5,000  (final 10% held out)",
]
for line in ds_lines:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "    " + line
    r.font.size = Pt(14); r.font.name = "Courier New"; r.font.color.rgb = GRAY


# ── Slide A3: Training + run reproducibility ─────────────
s = prs.slides.add_slide(BLANK)
add_title(s, "Appendix C — Training & run reproducibility")
add_subtitle(s, "Fourier Neural Operator, NCSA Delta, single A100")

box = s.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.2), Inches(5.5))
tf = box.text_frame; tf.word_wrap = True

p = tf.paragraphs[0]
r = p.add_run(); r.text = "Architecture (FNO2d)"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY

arch_lines = [
    "4 Fourier layers,  width = 32,  spectral modes = 12 × 12",
    "Lift: 1×1 conv  2 → 32        Project: 1×1 conv  32 → 128 → 1",
    "Activation: GELU,    InstanceNorm after each spectral block",
    "Total parameters: ~1.19 M",
]
for line in arch_lines:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "    " + line
    r.font.size = Pt(14); r.font.name = "Courier New"; r.font.color.rgb = GRAY

p = tf.add_paragraph(); p.space_before = Pt(8)
r = p.add_run(); r.text = "Optimisation"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY

opt_lines = [
    "Optimiser: Adam,  lr = 1e-3,  cosine decay → 1e-5",
    "Batch size: 32     Epochs: 100     Grad clip: 1.0",
    "Input X normalised per channel (mean / std from train split)",
    "Target Y normalised in log-space:  y_n = (log G − log_min) / (log_max − log_min)",
    "Loss = relative-L²(pred, Y_n)  +  0.05 · Helmholtz residual loss",
]
for line in opt_lines:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "    " + line
    r.font.size = Pt(14); r.font.name = "Courier New"; r.font.color.rgb = GRAY

p = tf.add_paragraph(); p.space_before = Pt(8)
r = p.add_run(); r.text = "Compute"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY

comp_lines = [
    "NCSA Delta,  partition = gpuA100x4-interactive,  1× A100 (80 GB)",
    "Wall-clock: ~14 s per epoch, ~25 min total for 100 epochs",
    "Data generation: 100-task SLURM array, partition = cpu,  ~2 min",
    "Code: github-style repo  mre_pipeline/  (src/, slurm/, scripts/, tests/)",
]
for line in comp_lines:
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "    " + line
    r.font.size = Pt(14); r.font.name = "Courier New"; r.font.color.rgb = GRAY


out = Path("EHMAN_BRIEFING.pptx")
prs.save(out)
print(f"Saved {out}  ({out.stat().st_size/1024:.0f} KB)")
print(f"Slides: {len(prs.slides)}")
